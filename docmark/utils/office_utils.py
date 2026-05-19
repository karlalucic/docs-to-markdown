"""
Utility functions for processing Office documents (DOCX, XLSX, PPTX).
"""

import logging
from dataclasses import dataclass
from pathlib import Path
from typing import List, Dict, Any, Tuple
import docx
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph
import openpyxl
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docmark.models.document_schema import ContentBlock, BlockType, TableContent, ImageContent

logger = logging.getLogger(__name__)


@dataclass
class DocxImagePayload:
    """DOCX body image payload for additive Vision descriptions."""
    sequence: int
    caption: str
    block_index: int
    mime_type: str
    image_bytes: bytes


class DOCXParser:
    """Parser for Microsoft Word (.docx) files."""

    _REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    @staticmethod
    def parse_document(docx_path: Path) -> List[ContentBlock]:
        """
        Parse DOCX document into content blocks.

        Compatibility wrapper around parse_document_with_images().
        """
        blocks, _ = DOCXParser.parse_document_with_images(docx_path)
        return blocks

    @staticmethod
    def parse_document_with_images(
        docx_path: Path,
    ) -> Tuple[List[ContentBlock], List[DocxImagePayload]]:
        """
        Parse DOCX document into content blocks.

        Args:
            docx_path: Path to DOCX file

        Returns:
            Tuple of (List of ContentBlock objects, List of image payloads)
        """
        doc = docx.Document(str(docx_path))
        blocks: List[ContentBlock] = []
        image_payloads: List[DocxImagePayload] = []
        sequence = 0

        for element in doc.element.body:
            # Check if it's a paragraph
            if element.tag.endswith("p"):
                para = Paragraph(element, doc)
                block = DOCXParser._parse_paragraph(para)
                if block:
                    blocks.append(block)

            # Check if it's a table
            elif element.tag.endswith("tbl"):
                table = DocxTable(element, doc)
                block = DOCXParser._parse_table(table)
                if block:
                    blocks.append(block)

            # Insert image blocks immediately after this body element's content blocks.
            # This preserves logical flow for chunking/indexing.
            for image_part in DOCXParser._extract_images_from_element(element, doc):
                sequence += 1
                caption = f"Embedded image {sequence}"
                blocks.append(
                    ContentBlock(
                        type=BlockType.IMAGE,
                        content=ImageContent(
                            description="",
                            caption=caption,
                            extraction_confidence=0.50,
                        ),
                    )
                )
                image_payloads.append(
                    DocxImagePayload(
                        sequence=sequence,
                        caption=caption,
                        block_index=len(blocks) - 1,
                        mime_type=image_part.content_type,
                        image_bytes=image_part.blob,
                    )
                )

        if image_payloads:
            logger.info(f"DOCX: Detected {len(image_payloads)} embedded image(s)")

        return blocks, image_payloads

    @staticmethod
    def _extract_images_from_element(element: Any, doc: docx.Document) -> List[Any]:
        """
        Extract image parts from one body element in XML order.
        Supports both inline and anchored drawings.
        """
        image_parts: List[Any] = []
        drawings = element.xpath('.//*[local-name()="drawing"]')
        for drawing in drawings:
            containers = drawing.xpath('./*[local-name()="inline" or local-name()="anchor"]')
            for container in containers:
                blips = container.xpath('.//*[local-name()="blip"]')
                for blip in blips:
                    rel_id = blip.get(f"{{{DOCXParser._REL_NS}}}embed")
                    if not rel_id:
                        continue

                    image_part = doc.part.related_parts.get(rel_id)
                    if image_part is None:
                        continue

                    image_bytes = getattr(image_part, "blob", None)
                    mime_type = getattr(image_part, "content_type", "")
                    if not image_bytes or not mime_type.startswith("image/"):
                        continue

                    image_parts.append(image_part)

        return image_parts

    @staticmethod
    def _parse_paragraph(para: Paragraph) -> ContentBlock:
        """Parse a paragraph into a ContentBlock."""
        text = para.text.strip()
        if not text:
            return None

        # Detect headings
        if para.style.name.startswith("Heading"):
            try:
                level = int(para.style.name.split()[-1])
            except (ValueError, IndexError):
                level = 1
            return ContentBlock(type=BlockType.HEADING, content=text, level=level)

        # Detect code
        if para.style.name in ["Code", "Source Code", "Preformatted"]:
            return ContentBlock(type=BlockType.CODE, content=text)

        # Regular text
        return ContentBlock(type=BlockType.TEXT, content=text)

    @staticmethod
    def _calculate_table_complexity(table: DocxTable) -> float:
        """
        Calculate table complexity and return appropriate confidence score.

        Analyzes merged cells, varying column counts, and nested tables.

        Args:
            table: DocxTable object

        Returns:
            float: 0.95 for simple tables, 0.75-0.85 for complex ones
        """
        complexity_score = 1.0

        # Check for merged cells (gridSpan for horizontal, vMerge for vertical)
        merged_cells = 0
        for row in table.rows:
            for cell in row.cells:
                try:
                    # Check horizontal merge (gridSpan)
                    if hasattr(cell._element, 'gridSpan') and cell._element.gridSpan is not None:
                        merged_cells += 1
                    # Check vertical merge (vMerge)
                    if hasattr(cell._element, 'vMerge') and cell._element.vMerge is not None:
                        merged_cells += 1
                except AttributeError:
                    pass

        if merged_cells > 0:
            complexity_score -= 0.1 * min(merged_cells / 5, 1.0)  # Max -0.1

        # Check for varying column counts (irregular structure)
        col_counts = [len(row.cells) for row in table.rows]
        if len(set(col_counts)) > 1:
            complexity_score -= 0.1

        # Check for nested tables (very complex)
        for row in table.rows:
            for cell in row.cells:
                if cell.tables:
                    complexity_score -= 0.15
                    break

        # Return final confidence (0.75 minimum to still be above Tier 2 threshold)
        return max(0.75, complexity_score * 0.95)

    @staticmethod
    def _table_to_html(table: DocxTable) -> str:
        """
        Convert complex table to HTML with colspan/rowspan for merged cells.

        Args:
            table: DocxTable object with merged cells

        Returns:
            HTML table string
        """
        html_lines = ["<table>"]

        for row_idx, row in enumerate(table.rows):
            html_lines.append("  <tr>")
            for cell in row.cells:
                # Detect column span
                colspan = 1
                try:
                    if hasattr(cell._element, 'gridSpan') and cell._element.gridSpan is not None:
                        colspan = int(cell._element.gridSpan)
                except (AttributeError, ValueError):
                    pass

                # Detect row span (simplified - vMerge is complex)
                rowspan = 1
                # Note: Full vMerge handling requires tracking merge starts/continues

                span_attrs = ""
                if colspan > 1:
                    span_attrs += f' colspan="{colspan}"'
                if rowspan > 1:
                    span_attrs += f' rowspan="{rowspan}"'

                tag = "th" if row_idx == 0 else "td"
                cell_text = cell.text.strip().replace('\n', '<br>')
                html_lines.append(f'    <{tag}{span_attrs}>{cell_text}</{tag}>')
            html_lines.append("  </tr>")

        html_lines.append("</table>")
        return "\n".join(html_lines)

    @staticmethod
    def _table_to_markdown(table: DocxTable) -> str:
        """
        Convert simple table to Markdown format.

        Args:
            table: DocxTable object without merged cells

        Returns:
            Markdown table string
        """
        # Process rows
        rows_data = []
        for row in table.rows:
            row_data = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows_data.append(row_data)

        if not rows_data:
            return ""

        # Assume first row is header
        header = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []

        # Build markdown table
        markdown_lines = []
        markdown_lines.append("| " + " | ".join(header) + " |")
        markdown_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        for row in data_rows:
            # Pad row if needed
            padded_row = row + [""] * (len(header) - len(row))
            markdown_lines.append("| " + " | ".join(padded_row[:len(header)]) + " |")

        return "\n".join(markdown_lines)

    @staticmethod
    def _parse_table(table: DocxTable) -> ContentBlock:
        """Parse a table with proper merged cell handling."""
        # Check if table has merged cells
        has_merged = any(
            (hasattr(cell._element, 'gridSpan') and cell._element.gridSpan is not None) or
            (hasattr(cell._element, 'vMerge') and cell._element.vMerge is not None)
            for row in table.rows
            for cell in row.cells
        )

        # Calculate complexity-based confidence
        confidence = DOCXParser._calculate_table_complexity(table)

        if has_merged:
            # Use HTML table for complex structure
            markdown = DOCXParser._table_to_html(table)
        else:
            # Use simple Markdown table
            markdown = DOCXParser._table_to_markdown(table)

        return ContentBlock(
            type=BlockType.TABLE,
            content=TableContent(
                markdown=markdown,
                caption=None,
                extraction_confidence=confidence,  # Now dynamic!
                enhanced_by_vision=False,
            ),
        )


class XLSXParser:
    """Parser for Microsoft Excel (.xlsx) files."""

    @staticmethod
    def parse_document(xlsx_path: Path) -> List[ContentBlock]:
        """
        Parse XLSX document into content blocks (one per sheet).

        Args:
            xlsx_path: Path to XLSX file

        Returns:
            List of ContentBlock objects (one table per sheet)
        """
        workbook = openpyxl.load_workbook(str(xlsx_path), data_only=True)
        blocks: List[ContentBlock] = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            block = XLSXParser._parse_sheet(sheet, sheet_name)
            if block:
                blocks.append(block)

        # Detect images and charts across all sheets
        image_count = 0
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            if hasattr(sheet, '_images') and sheet._images:
                for idx in range(len(sheet._images)):
                    image_count += 1
                    blocks.append(ContentBlock(
                        type=BlockType.IMAGE,
                        content=ImageContent(
                            description="",
                            caption=f"Image in sheet '{sheet_name}' ({idx + 1})",
                            extraction_confidence=0.50,
                        ),
                    ))
            if hasattr(sheet, '_charts') and sheet._charts:
                for idx in range(len(sheet._charts)):
                    image_count += 1
                    blocks.append(ContentBlock(
                        type=BlockType.IMAGE,
                        content=ImageContent(
                            description="",
                            caption=f"Chart in sheet '{sheet_name}' ({idx + 1})",
                            extraction_confidence=0.50,
                        ),
                    ))
        if image_count > 0:
            logger.info(f"XLSX: Detected {image_count} image(s)/chart(s)")

        return blocks

    @staticmethod
    def _calculate_sheet_complexity(sheet: Worksheet) -> float:
        """
        Calculate Excel sheet complexity and return appropriate confidence score.

        Analyzes merged cells and irregular structure.

        Args:
            sheet: openpyxl Worksheet object

        Returns:
            float: 0.95 for simple sheets, 0.75-0.85 for complex ones
        """
        complexity_score = 1.0

        # Check for merged cells
        merged_cell_count = len(sheet.merged_cells.ranges) if sheet.merged_cells else 0
        if merged_cell_count > 0:
            complexity_score -= 0.1 * min(merged_cell_count / 10, 1.0)  # Max -0.1

        # Check for varying column counts (irregular structure)
        col_counts = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                # Count non-empty cells
                non_empty = sum(1 for cell in row if cell is not None)
                col_counts.append(non_empty)

        if len(set(col_counts)) > 2:  # More than 2 different column counts
            complexity_score -= 0.05

        # Return final confidence (0.75 minimum)
        return max(0.75, complexity_score * 0.95)

    @staticmethod
    def _parse_sheet(sheet: Worksheet, sheet_name: str) -> ContentBlock:
        """Parse a worksheet into a ContentBlock with TableContent."""
        # Read all data
        data = []
        for row in sheet.iter_rows(values_only=True):
            # Skip completely empty rows
            if any(cell is not None for cell in row):
                data.append([str(cell) if cell is not None else "" for cell in row])

        if not data:
            return None

        # Calculate complexity-based confidence
        confidence = XLSXParser._calculate_sheet_complexity(sheet)

        # Convert to markdown table
        markdown_lines = []

        # First row as header
        header = data[0]
        data_rows = data[1:] if len(data) > 1 else []

        # Build markdown table
        markdown_lines.append("| " + " | ".join(header) + " |")
        markdown_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        for row in data_rows:
            # Pad row if needed
            padded_row = row + [""] * (len(header) - len(row))
            markdown_lines.append("| " + " | ".join(padded_row[:len(header)]) + " |")

        markdown = "\n".join(markdown_lines)

        return ContentBlock(
            type=BlockType.TABLE,
            content=TableContent(
                markdown=markdown,
                caption=f"Sheet: {sheet_name}",
                extraction_confidence=confidence,  # Now dynamic!
                enhanced_by_vision=False,
            ),
        )


class PPTXParser:
    """Parser for Microsoft PowerPoint (.pptx) files."""

    @staticmethod
    def parse_document(pptx_path: Path) -> List[Tuple[int, List[ContentBlock]]]:
        """
        Parse PPTX document into content blocks per slide.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            List of (slide_number, blocks) tuples (1-indexed)
        """
        prs = Presentation(str(pptx_path))
        slides_data: List[Tuple[int, List[ContentBlock]]] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            blocks: List[ContentBlock] = []

            for shape in slide.shapes:
                # Tables
                if shape.has_table:
                    block = PPTXParser._parse_table(shape.table)
                    if block:
                        blocks.append(block)

                # Text frames (headings and body text)
                elif shape.has_text_frame:
                    text_blocks = PPTXParser._parse_text_frame(shape.text_frame)
                    blocks.extend(text_blocks)

                # Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blocks.append(ContentBlock(
                        type=BlockType.IMAGE,
                        content=ImageContent(
                            description="",
                            caption=f"Slide {slide_num} image",
                            extraction_confidence=0.50,
                        ),
                    ))

                # Charts
                if shape.has_chart:
                    blocks.append(ContentBlock(
                        type=BlockType.IMAGE,
                        content=ImageContent(
                            description="",
                            caption=f"Slide {slide_num} chart",
                            extraction_confidence=0.50,
                        ),
                    ))

            slides_data.append((slide_num, blocks))

        return slides_data

    @staticmethod
    def _parse_text_frame(text_frame) -> List[ContentBlock]:
        """Parse a text frame into content blocks."""
        blocks: List[ContentBlock] = []

        for para in text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect headings by font size or level
            level = para.level
            if level == 0 and para.font and para.font.size and para.font.size >= Pt(20):
                blocks.append(ContentBlock(
                    type=BlockType.HEADING, content=text, level=1
                ))
            elif level == 0 and para.font and para.font.size and para.font.size >= Pt(16):
                blocks.append(ContentBlock(
                    type=BlockType.HEADING, content=text, level=2
                ))
            else:
                blocks.append(ContentBlock(type=BlockType.TEXT, content=text))

        return blocks

    @staticmethod
    def _parse_table(table) -> ContentBlock:
        """Parse a PPTX table into a ContentBlock."""
        rows_data = []
        has_merged = False

        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip().replace('\n', ' '))
                if cell.is_merge_origin or cell.span_height > 1 or cell.span_width > 1:
                    has_merged = True
            rows_data.append(row_data)

        if not rows_data:
            return None

        # Calculate confidence based on complexity
        confidence = 0.95
        if has_merged:
            confidence = 0.80

        # Build markdown table
        header = rows_data[0]
        data_rows = rows_data[1:] if len(rows_data) > 1 else []

        markdown_lines = []
        markdown_lines.append("| " + " | ".join(header) + " |")
        markdown_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        for row in data_rows:
            padded_row = row + [""] * (len(header) - len(row))
            markdown_lines.append("| " + " | ".join(padded_row[:len(header)]) + " |")

        return ContentBlock(
            type=BlockType.TABLE,
            content=TableContent(
                markdown="\n".join(markdown_lines),
                caption=None,
                extraction_confidence=confidence,
                enhanced_by_vision=False,
            ),
        )


def get_document_type(file_path: Path) -> str:
    """
    Determine document type from file extension.

    Args:
        file_path: Path to document file

    Returns:
        Document type: 'pdf', 'docx', 'xlsx', 'xls', 'pptx', or 'ppt'
    """
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return "pdf"
    elif suffix == ".docx":
        return "docx"
    elif suffix == ".xlsx":
        return "xlsx"
    elif suffix == ".xls":
        return "xls"
    elif suffix == ".pptx":
        return "pptx"
    elif suffix == ".ppt":
        return "ppt"
    else:
        raise ValueError(f"Unsupported file type: {suffix}")
